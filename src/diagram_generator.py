from diagrams import Diagram, Node, Edge
from pptx import Presentation
from pptx.util import Inches
import os


class ERDiagramGenerator:
    def __init__(self, entities, relationships):
        self.entities = entities
        self.relationships = relationships
        self.filename = None  # To store the base filename for reuse

    def generate(self, filename="optimized_er_diagram"):
        try:
            # Ensure the file ends with a single .png extension
            self.filename = filename.rstrip(".png")
            output_path = f"{self.filename}.png"  # Add .png extension for saving

            print(f"DEBUG: Generating diagram at: {output_path}")

            with Diagram(
                filename=self.filename,  # Pass the base name (without .png) to Diagram
                show=False,
                direction="LR",  # Left-to-right layout
            ):
                # Create graph nodes for each entity
                nodes = {
                    name: Node(label=self._entity_label(entity), shape="plaintext")
                    for name, entity in self.entities.items()
                }

                # Add edges for relationships
                for rel in self.relationships:
                    print(
                        f"DEBUG: Adding relationship: {rel.entity1} -- {rel.relation} --> {rel.entity2}"
                    )
                    nodes[rel.entity1] >> Edge(label=rel.relation) >> nodes[rel.entity2]

            # Check if the file was created
            print(f"DEBUG: Checking if file exists: {output_path}")
            if os.path.exists(output_path):  # Check the file with .png extension
                print(f"DEBUG: Diagram saved successfully at: {output_path}")
                return True
            else:
                print(f"ERROR: Diagram file not found after generation: {output_path}")
                return False
        except Exception as e:
            print(f"ERROR: Exception occurred during diagram generation: {e}")
            return False

    def export_to_pptx(self, pptx_filename="diagram.pptx"):
        """Exports the diagram to a PowerPoint file."""
        try:
            if not self.filename:
                raise ValueError("Diagram must be generated before exporting to PPTX.")

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

            # Add title to the slide
            title = slide.shapes.title
            title.text = "Entity-Relationship Diagram"

            # Add the ER diagram image to the slide
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            slide.shapes.add_picture(f"{self.filename}.png", left, top, width=width)

            # Save the PowerPoint file
            prs.save(pptx_filename)
            print(f"DEBUG: PowerPoint file saved at {pptx_filename}")
        except Exception as e:
            print(f"ERROR: Failed to export to PowerPoint: {e}")
            raise

    def export_to_dot(self, dot_filename="diagram.dot"):
        """Exports the diagram to a DOT file."""
        try:
            with open(dot_filename, "w") as f:
                f.write("digraph ERD {\n")
                f.write("rankdir=LR;\n")  # Left-to-right layout
                for entity_name, entity in self.entities.items():
                    f.write(f'"{entity_name}" [shape=plaintext, label=<{self._entity_label(entity)}>];\n')
                for rel in self.relationships:
                    f.write(f'"{rel.entity1}" -> "{rel.entity2}" [label="{rel.relation}"];\n')
                f.write("}\n")
            print(f"DEBUG: DOT file saved at {dot_filename}")
        except Exception as e:
            print(f"ERROR: Failed to export DOT file: {e}")
            raise

    def _entity_label(self, entity):
        """Creates an HTML-style table for the entity with a specific format."""
        label = (
            f'<<TABLE BORDER="1" CELLBORDER="1" CELLSPACING="0" CELLPADDING="4">'
        )
        label += f'<TR><TD COLSPAN="2" BGCOLOR="lightgrey" ALIGN="CENTER"><B>{entity.name}</B></TD></TR>'

        for attr in entity.attributes:
            pk_indicator = "<B>PK </B>" if attr.is_primary else ""
            label += f'<TR>'
            label += (
                f'<TD ALIGN="LEFT">{pk_indicator}<B>{attr.name}</B></TD>'
            )  # Bold attribute name
            label += (
                f'<TD ALIGN="LEFT">{self._get_data_type(attr)}</TD>'
            )  # Data type
            label += f'</TR>'

        label += "</TABLE>>"
        return label

    def _get_data_type(self, attr):
        # Example data type mapping; customize as needed
        if "id" in attr.name.lower():
            return "integer"
        elif "name" in attr.name.lower():
            return "string"
        elif "price" in attr.name.lower():
            return "float"
        elif "bool" in attr.name.lower() or "eligible" in attr.name.lower():
            return "bool"
        else:
            return "string"
